from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Color Palette (Premium Dark Theme)
    BG_COLOR = RGBColor(11, 15, 25)        # #0B0F19 - Deep slate
    TITLE_COLOR = RGBColor(165, 180, 252) # #A5B4FC - Light indigo
    GOLD_COLOR = RGBColor(251, 191, 36)   # #FBBF24 - Golden amber
    TEXT_COLOR = RGBColor(226, 232, 240)   # #E2E8F0 - Soft white
    ACCENT_COLOR = RGBColor(108, 99, 255) # #6C63FF - Neon purple
    BOX_BG = RGBColor(20, 27, 45)         # #141B2D - Dark container card
    BORDER_COLOR = RGBColor(51, 65, 85)   # #334155 - Slate border

    def set_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR

    def add_slide_frame(title_text, category_text=""):
        slide = prs.slides.add_slide(prs.slide_layouts[6]) # blank layout
        set_background(slide)
        
        # Category Tag
        if category_text:
            cat_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.4), Inches(11.83), Inches(0.4))
            p = cat_box.text_frame.paragraphs[0]
            p.text = category_text.upper()
            p.font.name = 'Arial'
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.color.rgb = GOLD_COLOR
            
        # Slide Title
        title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.7), Inches(11.83), Inches(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = 'Arial'
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = TITLE_COLOR
        
        return slide

    def add_text_box(slide, text_lines, left, top, width, height, font_size=14, bold_first=False):
        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = True
        for i, line in enumerate(text_lines):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line
            p.font.name = 'Arial'
            p.font.size = Pt(font_size)
            p.font.color.rgb = TEXT_COLOR
            p.space_after = Pt(8)
            if bold_first and i == 0:
                p.font.bold = True
                p.font.size = Pt(font_size + 2)
                p.font.color.rgb = GOLD_COLOR
        return box

    def add_image_or_placeholder(slide, img_path, left, top, width, height, label):
        if os.path.exists(img_path):
            slide.shapes.add_picture(img_path, left, top, width=width, height=height)
        else:
            # Add a styled placeholder shape
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = BOX_BG
            shape.line.color.rgb = BORDER_COLOR
            shape.line.width = Pt(1.5)
            
            tf = shape.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = f"[ Live screenshot of {label} will be inserted here ]\n\nRun 'python capture_screenshots.py' to automatically populate this slide."
            p.alignment = PP_ALIGN.CENTER
            p.font.name = 'Arial'
            p.font.size = Pt(12)
            p.font.color.rgb = muted_color = RGBColor(148, 163, 184)

    # ----------------------------------------------------
    # SLIDE 1: Title Slide
    # ----------------------------------------------------
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide1)
    
    # Styled Accent Shape (Background Glow Effect)
    glow = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.4), Inches(7.5))
    glow.fill.solid()
    glow.fill.fore_color.rgb = ACCENT_COLOR
    glow.line.fill.background()

    # Title Text Frame
    title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.3), Inches(3.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    
    p1 = tf.paragraphs[0]
    p1.text = "AdaptLearn"
    p1.font.name = 'Arial'
    p1.font.size = Pt(64)
    p1.font.bold = True
    p1.font.color.rgb = GOLD_COLOR
    p1.space_after = Pt(5)
    
    p2 = tf.add_paragraph()
    p2.text = "An AI-Powered Adaptive Learning & Evaluation Platform"
    p2.font.name = 'Arial'
    p2.font.size = Pt(28)
    p2.font.color.rgb = TITLE_COLOR
    p2.space_after = Pt(30)
    
    p3 = tf.add_paragraph()
    p3.text = "Django Web Development Project Submission"
    p3.font.name = 'Arial'
    p3.font.size = Pt(14)
    p3.font.color.rgb = TEXT_COLOR
    
    p4 = tf.add_paragraph()
    p4.text = "Presented by: Student Name"
    p4.font.name = 'Arial'
    p4.font.size = Pt(14)
    p4.font.color.rgb = TEXT_COLOR

    # ----------------------------------------------------
    # SLIDE 2: Project Overview & Motivation
    # ----------------------------------------------------
    slide2 = add_slide_frame("Project Overview & Motivation", "Introduction")
    overview_text = [
        "Traditional e-learning platforms are static and treat all students equally, regardless of their prior knowledge, learning styles, or speed of assimilation.",
        "AdaptLearn solves this by building a dynamic learning system that tailors recommendations, evaluates student progress, and provides instant, personalized guidance.",
        "• Personalization: Dynamic adjustments based on student skill levels (Beginner, Intermediate, Advanced).",
        "• AI Integration: Real-time tutoring via AI Mentor Chatbot and automatic grading feedback.",
        "• Engagement: Gamification with XP points, progress badges, and ranking leaderboard."
    ]
    add_text_box(slide2, overview_text, Inches(0.75), Inches(1.8), Inches(11.8), Inches(4.5), font_size=15)

    # ----------------------------------------------------
    # SLIDE 3: Key Features & Architecture
    # ----------------------------------------------------
    slide3 = add_slide_frame("Key Architecture & Features", "Platform Features")
    
    features_col1 = [
        "Core Features",
        "• Adaptive Recommendations: Custom course suggestions dynamically tailored to student's profile skill level.",
        "• AI Mentor (AdaptBot): Seamless conversational AI mentor to clarify doubt & code snippets.",
        "• AI Exam Evaluation: Automatic grading and contextual AI feedback on exam submissions.",
        "• Gamified System: XP points awarded upon material completion, exams passed, and course completions."
    ]
    features_col2 = [
        "System Architecture",
        "• Backend Framework: Django (MVT architecture) ensuring high performance, security, and SQLite DB integration.",
        "• Premium Frontend: Vanilla CSS with custom glassmorphism components, dark-mode elements, and Chart.js dashboards.",
        "• Document Generation: Dynamic generation of PDF Progress Reports and verifiable Certificates of Completion."
    ]
    add_text_box(slide3, features_col1, Inches(0.75), Inches(1.8), Inches(5.6), Inches(4.5), font_size=14, bold_first=True)
    add_text_box(slide3, features_col2, Inches(6.9), Inches(1.8), Inches(5.6), Inches(4.5), font_size=14, bold_first=True)

    # ----------------------------------------------------
    # SLIDE 4: Database Schema & Key Models
    # ----------------------------------------------------
    slide4 = add_slide_frame("Database Design & Core Models", "Data Structure")
    
    models_col1 = [
        "User & Profile Models",
        "• UserProfile: Extends base User, handles role (student/instructor/admin), total XP points, and skill level.",
        "• Enrollment: Manages course registrations, tracking progress %, completion timestamps, and certificate flags.",
        "• StudentProgress: Keeps track of individual module materials completed by students."
    ]
    models_col2 = [
        "Course & Evaluation Models",
        "• Course / Module / LearningMaterial: Represents the hierarchy of structured study content.",
        "• Exam / Question: Supports multiple-choice or true/false formats with custom marks and difficulties.",
        "• ExamResult / StudentAnswer: Stores user answers, computed percentages, pass/fail status, and AI feedback."
    ]
    add_text_box(slide4, models_col1, Inches(0.75), Inches(1.8), Inches(5.6), Inches(4.5), font_size=14, bold_first=True)
    add_text_box(slide4, models_col2, Inches(6.9), Inches(1.8), Inches(5.6), Inches(4.5), font_size=14, bold_first=True)

    # ----------------------------------------------------
    # SLIDE 5: Student Dashboard Output
    # ----------------------------------------------------
    slide5 = add_slide_frame("Student Dashboard Output", "User Interface")
    dashboard_desc = [
        "Key Metrics & Navigation",
        "• Enrolled Courses count, Completed courses metric (correctly tracking 100% complete courses), and Total XP points.",
        "• Quick access to course list, exam listings, notifications, and adaptive recommendations.",
        "• Status flags showing 'Completed' with a certificate link for graduated courses, instead of standard continue buttons."
    ]
    add_text_box(slide5, dashboard_desc, Inches(0.75), Inches(1.8), Inches(5.0), Inches(4.5), font_size=14, bold_first=True)
    add_image_or_placeholder(slide5, "dashboard.png", Inches(6.2), Inches(1.8), Inches(6.5), Inches(4.5), "Student Dashboard")

    # ----------------------------------------------------
    # SLIDE 6: Course Content & Modules Output
    # ----------------------------------------------------
    slide6 = add_slide_frame("Course Material & Detail View", "User Interface")
    content_desc = [
        "Course Navigation",
        "• Displays course syllabus, total duration, difficulty badge, and enrolled students count.",
        "• Collapse-accordion modules list containing videos, articles, and quizzes.",
        "• Gold 'Claim Certificate' action button dynamically appears in the sidebar as soon as the student hits 100% completion.",
        "• 'Restart Course' button resets progress, clearing progress items and allowing a clean start."
    ]
    add_text_box(slide6, content_desc, Inches(0.75), Inches(1.8), Inches(5.0), Inches(4.5), font_size=14, bold_first=True)
    add_image_or_placeholder(slide6, "course_detail.png", Inches(6.2), Inches(1.8), Inches(6.5), Inches(4.5), "Course Detail & Modules")

    # ----------------------------------------------------
    # SLIDE 7: AI Mentor Chatbot (AdaptBot) Output
    # ----------------------------------------------------
    slide7 = add_slide_frame("AI Mentor Chatbot (AdaptBot)", "AI Integration")
    chatbot_desc = [
        "AdaptBot Interface",
        "• Integrated Anthropic Claude API to provide conversational, friendly, and helpful tutoring support.",
        "• Preserves student's chat history in database, providing contextual study support.",
        "• Able to clarify code snippets, break down concepts, and suggest pomodoro study goals."
    ]
    add_text_box(slide7, chatbot_desc, Inches(0.75), Inches(1.8), Inches(5.0), Inches(4.5), font_size=14, bold_first=True)
    add_image_or_placeholder(slide7, "chatbot.png", Inches(6.2), Inches(1.8), Inches(6.5), Inches(4.5), "AdaptBot Chat Interface")

    # ----------------------------------------------------
    # SLIDE 8: Student Performance Analytics
    # ----------------------------------------------------
    slide8 = add_slide_frame("Performance Analytics Dashboard", "User Interface")
    analytics_desc = [
        "Metrics Visualization",
        "• Exam Score Trend: Multi-line Chart.js chart comparing scores against the pass boundary.",
        "• Subject Distribution: Doughnut chart visualizing attempts categorised by type (Module, Final, Practice).",
        "• Activity Map: 12-week grid heatmap tracking daily learning actions.",
        "• Radar Distribution: Visualizes score percentages grouped by difficulty (Easy, Medium, Hard)."
    ]
    add_text_box(slide8, analytics_desc, Inches(0.75), Inches(1.8), Inches(5.0), Inches(4.5), font_size=14, bold_first=True)
    add_image_or_placeholder(slide8, "analytics.png", Inches(6.2), Inches(1.8), Inches(6.5), Inches(4.5), "Analytics & Chart Dashboard")

    # ----------------------------------------------------
    # SLIDE 9: Certificate of Completion
    # ----------------------------------------------------
    slide9 = add_slide_frame("Dynamic Completion Certificate", "System Output")
    certificate_desc = [
        "Verifiable Certificates",
        "• Generated dynamically using HTML/CSS print layers when course completion reaches 100%.",
        "• Features the student's highest passing exam score, completion date, and verification certificate ID.",
        "• Includes built-in print control stylesheets to trigger clean physical prints or PDF exports directly from browser."
    ]
    add_text_box(slide9, certificate_desc, Inches(0.75), Inches(1.8), Inches(5.0), Inches(4.5), font_size=14, bold_first=True)
    add_image_or_placeholder(slide9, "certificate.png", Inches(6.2), Inches(1.8), Inches(6.5), Inches(4.5), "Verifiable Certificate")

    # ----------------------------------------------------
    # SLIDE 10: Dynamic PDF Progress Report
    # ----------------------------------------------------
    slide10 = add_slide_frame("Structured PDF Progress Report", "System Output")
    pdf_desc = [
        "PDF Report Generation",
        "• Exports comprehensive student summary in a beautifully formatted PDF document.",
        "• Utilizes ReportLab flowables to render structured tables, header brand, and student metadata.",
        "• Generates list of registered courses with progress metrics, alongside full historical exam logs with color-coded status badges."
    ]
    add_text_box(slide10, pdf_desc, Inches(0.75), Inches(1.8), Inches(5.0), Inches(4.5), font_size=14, bold_first=True)
    add_image_or_placeholder(slide10, "pdf_report.png", Inches(6.2), Inches(1.8), Inches(6.5), Inches(4.5), "Exported PDF Report Page")

    # ----------------------------------------------------
    # SLIDE 11: Conclusion & Future Scope
    # ----------------------------------------------------
    slide11 = add_slide_frame("Conclusion & Future Scope", "Project Summary")
    conclusion_text = [
        "Project Accomplishments",
        "• Successfully developed an AI-driven e-learning ecosystem utilizing Django.",
        "• Solved user progression tracking, exam evaluation, and visual analytics.",
        "• Implemented functional gamification with verified XP accumulation.",
        "",
        "Future Enhancements",
        "• Multi-media adaptive materials generation (using generative AI on the fly).",
        "• Integration with standard LMS APIs (SCORM / LTI compatibility).",
        "• Mobile application companion using React Native / Flutter."
    ]
    add_text_box(slide11, conclusion_text, Inches(0.75), Inches(1.8), Inches(11.8), Inches(4.5), font_size=15)

    prs.save("AdaptLearn_Project_Presentation.pptx")
    print("Presentation saved successfully as AdaptLearn_Project_Presentation.pptx!")

if __name__ == "__main__":
    create_presentation()
